from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[2]
SOURCE_DIR = ROOT / 'converters' / 'factory_tour'
HTML_PATH = SOURCE_DIR / '工場見学.html'
OUTPUT_PATH = ROOT / 'documentation' / 'presentations' / 'FactoryTour_Presentation.pptx'
REFERENCE_FONT = 'Yu Gothic'

with HTML_PATH.open('r', encoding='utf-8') as f:
    soup = BeautifulSoup(f, 'lxml')

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

left_margin = Inches(0.4)
right_margin = Inches(0.4)
top_margin = Inches(0.3)
content_width = prs.slide_width - left_margin - right_margin

# Background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
background.line.width = 0

# Header bar
header = soup.select_one('.header')
header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.0))
header_shape.fill.solid()
header_shape.fill.fore_color.rgb = RGBColor(0x2C, 0x3E, 0x50)
header_shape.line.width = 0

header_text_frame = header_shape.text_frame
header_text_frame.clear()
if header:
    title = header.find('h1').get_text(strip=True) if header.find('h1') else ''
    subtitle = header.select_one('.subtitle').get_text(strip=True) if header.select_one('.subtitle') else ''
else:
    title = ''
    subtitle = ''

p = header_text_frame.paragraphs[0]
p.text = title
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.font.name = REFERENCE_FONT

if subtitle:
    p = header_text_frame.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    p.font.name = REFERENCE_FONT

# Section columns
column_gap = Inches(0.3)
left_width = (content_width - column_gap) * 0.52
right_width = content_width - column_gap - left_width
left_x = left_margin
right_x = left_margin + left_width + column_gap
current_y = top_margin + Inches(1.05)

# Left intro box
intro = soup.select_one('.intro')
if intro:
    intro_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, current_y, left_width, Inches(2.1))
    intro_shape.fill.solid()
    intro_shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    intro_shape.line.width = Pt(0)
    intro_shape.shadow.inherit = False
    intro_text = intro_shape.text_frame
    intro_text.clear()

    h2 = intro.find('h2')
    p = intro_text.paragraphs[0]
    p.text = h2.get_text(strip=True) if h2 else '概要'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p.font.name = REFERENCE_FONT

    body_paragraph = intro_text.add_paragraph()
    body_paragraph.text = intro.find('p').get_text(strip=True) if intro.find('p') else ''
    body_paragraph.font.size = Pt(12)
    body_paragraph.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    body_paragraph.font.name = REFERENCE_FONT
    body_paragraph.space_before = Pt(6)

    current_y += intro_shape.height + Inches(0.2)

# Comparison table
comparison_table = soup.select_one('.comparison-table table')
if comparison_table:
    rows = comparison_table.find_all('tr')
    cols = rows[0].find_all(['th', 'td'])
    table_height = Inches(3.6)
    table_shape = slide.shapes.add_table(len(rows), len(cols), left_x, current_y, left_width, table_height)
    table = table_shape.table

    for i, row in enumerate(rows):
        cells = row.find_all(['th', 'td'])
        for j, cell in enumerate(cells):
            cell_text = ' '.join(cell.stripped_strings)
            table.cell(i, j).text = cell_text
            for paragraph in table.cell(i, j).text_frame.paragraphs:
                paragraph.font.name = REFERENCE_FONT
                paragraph.font.size = Pt(10 if i == 0 else 9)
                paragraph.space_after = 0
                paragraph.space_before = 0
                if i == 0:
                    paragraph.font.bold = True
    for j in range(len(cols)):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x34, 0x98, 0xDB)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            paragraph.alignment = PP_ALIGN.CENTER
    current_y += table_height + Inches(0.2)

# Bottom key insights block below table if space, else adjust
def add_key_insights(slide, x, y, width):
    bottom = soup.select_one('.key-insights')
    if not bottom:
        return
    height = Inches(1.5)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.color.rgb = RGBColor(0x34, 0x98, 0xDB)
    box.line.width = Pt(1)

    tf = box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True
    tf.clear()

    title = bottom.find('h3').get_text(strip=True)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
    p.font.name = REFERENCE_FONT

    items = bottom.select('.insight-item')
    for item in items:
        icon = item.select_one('.insight-icon').get_text(strip=True)
        text = item.find_all('div')[-1].get_text(strip=True)
        paragraph = tf.add_paragraph()
        paragraph.text = f"{icon}. {text}"
        paragraph.font.size = Pt(12)
        paragraph.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
        paragraph.font.name = REFERENCE_FONT
        paragraph.level = 1

    return height

# Right column content
right_y = top_margin + Inches(1.05)
strategy_section = soup.select_one('.strategy-section')
if strategy_section:
    heading = strategy_section.find('h2').get_text(strip=True)
    section_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, right_y, right_width, Inches(3.5))
    section_box.fill.solid()
    section_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    section_box.line.width = 0
    tf = section_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p.font.name = REFERENCE_FONT

    for item in strategy_section.select('.strategy-item'):
        subtitle = item.find('h3').get_text(strip=True)
        body = item.find('p').get_text(strip=True)
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.bold = True
        sp.font.color.rgb = RGBColor(0x34, 0x98, 0xDB)
        sp.font.name = REFERENCE_FONT
        sp.space_before = Pt(8)

        bp = tf.add_paragraph()
        bp.text = body
        bp.font.size = Pt(12)
        bp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        bp.font.name = REFERENCE_FONT
        bp.level = 1
    right_y += section_box.height + Inches(0.25)

conclusion = soup.select_one('.conclusion')
if conclusion:
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, right_y, right_width, Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    box.line.width = 0
    tf = box.text_frame
    tf.clear()

    heading = conclusion.find('h2').get_text(strip=True)
    p = tf.paragraphs[0]
    p.text = heading
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    p.font.name = REFERENCE_FONT

    for para in conclusion.find_all('p'):
        bp = tf.add_paragraph()
        bp.text = para.get_text(strip=True)
        bp.font.size = Pt(12)
        bp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        bp.font.name = REFERENCE_FONT
        bp.level = 1
        bp.space_before = Pt(4)
    right_y += box.height + Inches(0.25)

add_key_insights(slide, left_x, prs.slide_height - Inches(1.7), left_width)

OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUTPUT_PATH)
print(f'Presentation saved to {OUTPUT_PATH}')
